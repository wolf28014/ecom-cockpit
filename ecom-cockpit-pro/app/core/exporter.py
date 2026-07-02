"""
报表导出服务 - Excel/PDF/Word/PPT
"""
from datetime import date, datetime, timedelta
from pathlib import Path
from typing import List, Optional
import os

from sqlalchemy.orm import Session

from app.config import EXPORT_DIR
from app.database.models import DailyRecord, Store
from app.core.analytics import AnalyticsService


class ExportService:
    """报表导出服务"""

    def __init__(self, session: Session):
        self.session = session
        self.analytics = AnalyticsService(session)

    def _ensure_dir(self, file_path: Path) -> Path:
        file_path.parent.mkdir(parents=True, exist_ok=True)
        return file_path

    def _generate_filename(self, prefix: str, ext: str, store_id: int = 0) -> Path:
        today = date.today().isoformat()
        store_suffix = f"_store{store_id}" if store_id else "_all"
        filename = f"{prefix}_{today}{store_suffix}.{ext}"
        return self._ensure_dir(EXPORT_DIR / filename)

    # ---------- Excel ----------
    def export_daily_excel(self, store_id: int = 0, days: int = 30) -> Path:
        """导出每日数据 Excel"""
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        records = self.analytics.get_daily_records(store_id, days)
        store_name = self._get_store_name(store_id)

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "每日经营数据"

        # 标题
        ws.merge_cells("A1:L1")
        ws["A1"] = f"{store_name} - 每日经营数据（最近 {days} 天）"
        ws["A1"].font = Font(size=16, bold=True, color="0071E3")
        ws["A1"].alignment = Alignment(horizontal="center", vertical="center")

        # 表头
        headers = ["日期", "销售额", "订单数", "客单价", "退款金额", "退款率",
                   "推广费", "成本", "毛利润", "净利润", "利润率", "ROI"]
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=3, column=col, value=h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="0071E3")
            cell.alignment = Alignment(horizontal="center")

        # 数据
        for row_idx, r in enumerate(records, 4):
            ws.cell(row=row_idx, column=1, value=r.record_date.isoformat())
            ws.cell(row=row_idx, column=2, value=r.sales_amount).number_format = "¥#,##0.00"
            ws.cell(row=row_idx, column=3, value=r.order_count)
            ws.cell(row=row_idx, column=4, value=r.avg_order_value).number_format = "¥#,##0.00"
            ws.cell(row=row_idx, column=5, value=r.refund_amount).number_format = "¥#,##0.00"
            ws.cell(row=row_idx, column=6, value=r.refund_rate).number_format = "0.00%"
            ws.cell(row=row_idx, column=7, value=r.promotion_total).number_format = "¥#,##0.00"
            ws.cell(row=row_idx, column=8, value=r.cost_total).number_format = "¥#,##0.00"
            ws.cell(row=row_idx, column=9, value=r.gross_profit).number_format = "¥#,##0.00"
            ws.cell(row=row_idx, column=10, value=r.net_profit).number_format = "¥#,##0.00"
            ws.cell(row=row_idx, column=11, value=r.profit_rate).number_format = "0.00%"
            ws.cell(row=row_idx, column=12, value=r.roi).number_format = "0.00"

        # 列宽
        for col in range(1, 13):
            ws.column_dimensions[get_column_letter(col)].width = 14

        # 保存
        file_path = self._generate_filename("daily_report", "xlsx", store_id)
        wb.save(file_path)
        return file_path

    def export_summary_excel(self, store_id: int = 0) -> Path:
        """导出经营汇总 Excel"""
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment

        today = self.analytics.get_today_summary(store_id)
        week = self.analytics.get_week_summary(store_id)
        month = self.analytics.get_month_summary(store_id)
        year = self.analytics.get_year_summary(store_id)
        store_name = self._get_store_name(store_id)

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "经营汇总"

        ws.merge_cells("A1:D1")
        ws["A1"] = f"{store_name} - 经营汇总报告（{date.today().isoformat()}）"
        ws["A1"].font = Font(size=16, bold=True, color="0071E3")
        ws["A1"].alignment = Alignment(horizontal="center")

        # 表头
        headers = ["指标", "今日", "本周", "本月", "本年"]
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=3, column=col, value=h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="0071E3")
            cell.alignment = Alignment(horizontal="center")

        rows = [
            ("销售额", today.sales_amount, week.sales_amount, month.sales_amount, year.sales_amount),
            ("订单数", today.order_count, week.order_count, month.order_count, year.order_count),
            ("净利润", today.net_profit, week.net_profit, month.net_profit, year.net_profit),
            ("利润率", today.profit_rate, week.profit_rate, month.profit_rate, year.profit_rate),
            ("ROI", today.roi, week.roi, month.roi, year.roi),
            ("客单价", today.avg_order_value, week.avg_order_value, month.avg_order_value, year.avg_order_value),
            ("退款率", today.refund_rate, week.refund_rate, month.refund_rate, year.refund_rate),
            ("推广费率", today.promotion_rate, week.promotion_rate, month.promotion_rate, year.promotion_rate),
        ]
        for row_idx, (label, *vals) in enumerate(rows, 4):
            ws.cell(row=row_idx, column=1, value=label).font = Font(bold=True)
            for col, v in enumerate(vals, 2):
                cell = ws.cell(row=row_idx, column=col, value=v)
                if "率" in label:
                    cell.number_format = "0.00%"
                elif label in ("订单数",):
                    cell.number_format = "0"
                else:
                    cell.number_format = "¥#,##0.00"

        # 列宽
        for col, w in enumerate([14, 16, 16, 16, 16], 1):
            ws.column_dimensions[chr(64 + col)].width = w

        file_path = self._generate_filename("summary_report", "xlsx", store_id)
        wb.save(file_path)
        return file_path

    # ---------- PDF ----------
    def export_pdf(self, store_id: int = 0, period: str = "month") -> Path:
        """导出 PDF 报告"""
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        )
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        # 注册中文字体 + 字体族（ReportLab 需要 family 才能识别 bold/italic）
        from reportlab.pdfbase.pdfmetrics import registerFontFamily
        from pathlib import Path as _P

        # Noto Serif SC
        noto_serif_regular = "/usr/share/fonts/truetype/noto-serif-sc/NotoSerifSC-Regular.ttf"
        noto_serif_bold = "/usr/share/fonts/truetype/noto-serif-sc/NotoSerifSC-Bold.ttf"
        if _P(noto_serif_regular).exists():
            try:
                pdfmetrics.registerFont(TTFont("NotoSerifSC", noto_serif_regular))
                if _P(noto_serif_bold).exists():
                    pdfmetrics.registerFont(TTFont("NotoSerifSC-Bold", noto_serif_bold))
                else:
                    pdfmetrics.registerFont(TTFont("NotoSerifSC-Bold", noto_serif_regular))
                registerFontFamily(
                    "NotoSerifSC",
                    normal="NotoSerifSC",
                    bold="NotoSerifSC-Bold",
                    italic="NotoSerifSC",
                    boldItalic="NotoSerifSC-Bold",
                )
            except Exception:
                pass

        # Noto Sans SC（备用）
        noto_sans_regular = "/usr/share/fonts/truetype/chinese/NotoSansSC-Regular.ttf"
        if _P(noto_sans_regular).exists():
            try:
                pdfmetrics.registerFont(TTFont("NotoSansSC", noto_sans_regular))
                registerFontFamily(
                    "NotoSansSC",
                    normal="NotoSansSC",
                    bold="NotoSansSC",
                    italic="NotoSansSC",
                    boldItalic="NotoSansSC",
                )
            except Exception:
                pass

        # 数据
        if period == "today":
            summary = self.analytics.get_today_summary(store_id)
            period_label = "今日"
        elif period == "week":
            summary = self.analytics.get_week_summary(store_id)
            period_label = "本周"
        elif period == "year":
            summary = self.analytics.get_year_summary(store_id)
            period_label = "本年"
        else:
            summary = self.analytics.get_month_summary(store_id)
            period_label = "本月"

        store_name = self._get_store_name(store_id)
        file_path = self._generate_filename(f"{period}_report", "pdf", store_id)

        doc = SimpleDocTemplate(
            str(file_path), pagesize=A4,
            leftMargin=20*mm, rightMargin=20*mm,
            topMargin=20*mm, bottomMargin=20*mm
        )

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "CustomTitle", parent=styles["Heading1"],
            fontName="NotoSerifSC", fontSize=20,
            textColor=colors.HexColor("#0071E3"),
            alignment=1, spaceAfter=20
        )
        body_style = ParagraphStyle(
            "CustomBody", parent=styles["Normal"],
            fontName="NotoSerifSC", fontSize=10,
            leading=18, spaceAfter=10
        )
        heading_style = ParagraphStyle(
            "CustomHeading", parent=styles["Heading2"],
            fontName="NotoSerifSC", fontSize=14,
            textColor=colors.HexColor("#1D1D1F"),
            spaceAfter=10, spaceBefore=15
        )

        elements = []
        elements.append(Paragraph(f"{store_name}", title_style))
        elements.append(Paragraph(f"{period_label}经营报告 - {date.today().isoformat()}", body_style))
        elements.append(Spacer(1, 10))

        # 核心数据表
        elements.append(Paragraph("一、核心数据", heading_style))
        data = [
            ["指标", "数值"],
            ["销售额", f"¥{summary.sales_amount:,.2f}"],
            ["订单数", f"{summary.order_count:,}"],
            ["净利润", f"¥{summary.net_profit:,.2f}"],
            ["利润率", f"{summary.profit_rate*100:.2f}%"],
            ["ROI", f"{summary.roi:.2f}"],
            ["客单价", f"¥{summary.avg_order_value:.2f}"],
            ["推广费率", f"{summary.promotion_rate*100:.2f}%"],
            ["退款率", f"{summary.refund_rate*100:.2f}%"],
        ]
        table = Table(data, colWidths=[60*mm, 80*mm])
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0071E3")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, -1), "NotoSerifSC"),
            ("FONTSIZE", (0, 0), (-1, -1), 11),
            ("ALIGN", (0, 0), (-1, -1), "LEFT"),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#E5E5EA")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F5F5F7")]),
            ("TOPPADDING", (0, 0), (-1, -1), 8),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ]))
        elements.append(table)
        elements.append(Spacer(1, 15))

        # 分析
        elements.append(Paragraph("二、经营分析", heading_style))
        analysis_text = f"""
        {period_label}共完成销售额 <b>¥{summary.sales_amount:,.2f}</b>，净利润 <b>¥{summary.net_profit:,.2f}</b>，
        整体利润率 <b>{summary.profit_rate*100:.2f}%</b>，ROI 为 <b>{summary.roi:.2f}</b>。
        期间共完成订单 <b>{summary.order_count:,}</b> 笔，客单价 <b>¥{summary.avg_order_value:.2f}</b>，
        推广费率 <b>{summary.promotion_rate*100:.2f}%</b>，退款率 <b>{summary.refund_rate*100:.2f}%</b>。
        """
        elements.append(Paragraph(analysis_text, body_style))

        elements.append(Paragraph("三、说明", heading_style))
        elements.append(Paragraph(
            "本报告由电商经营驾驶舱 Pro 自动生成。如需详细分析，请使用 AI 经营中心生成智能分析报告。",
            body_style
        ))

        doc.build(elements)
        return file_path

    # ---------- Word ----------
    def export_word(self, store_id: int = 0, period: str = "month") -> Path:
        """导出 Word 报告"""
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        if period == "today":
            summary = self.analytics.get_today_summary(store_id)
            period_label = "今日"
        elif period == "week":
            summary = self.analytics.get_week_summary(store_id)
            period_label = "本周"
        elif period == "year":
            summary = self.analytics.get_year_summary(store_id)
            period_label = "本年"
        else:
            summary = self.analytics.get_month_summary(store_id)
            period_label = "本月"

        store_name = self._get_store_name(store_id)
        doc = Document()

        # 标题
        title = doc.add_heading(f"{store_name}", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title.runs:
            run.font.color.rgb = RGBColor(0x00, 0x71, 0xE3)

        subtitle = doc.add_paragraph(f"{period_label}经营报告 - {date.today().isoformat()}")
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph("")

        # 核心数据
        doc.add_heading("一、核心数据", level=1)
        table = doc.add_table(rows=9, cols=2)
        table.style = "Light Grid Accent 1"
        rows_data = [
            ("指标", "数值"),
            ("销售额", f"¥{summary.sales_amount:,.2f}"),
            ("订单数", f"{summary.order_count:,}"),
            ("净利润", f"¥{summary.net_profit:,.2f}"),
            ("利润率", f"{summary.profit_rate*100:.2f}%"),
            ("ROI", f"{summary.roi:.2f}"),
            ("客单价", f"¥{summary.avg_order_value:.2f}"),
            ("推广费率", f"{summary.promotion_rate*100:.2f}%"),
            ("退款率", f"{summary.refund_rate*100:.2f}%"),
        ]
        for row_idx, (label, value) in enumerate(rows_data):
            cells = table.rows[row_idx].cells
            cells[0].text = label
            cells[1].text = value

        # 经营分析
        doc.add_heading("二、经营分析", level=1)
        p = doc.add_paragraph()
        p.add_run(f"{period_label}共完成销售额 ").font.size = Pt(11)
        p.add_run(f"¥{summary.sales_amount:,.2f}").bold = True
        p.add_run("，净利润 ")
        p.add_run(f"¥{summary.net_profit:,.2f}").bold = True
        p.add_run(f"，整体利润率 ")
        p.add_run(f"{summary.profit_rate*100:.2f}%").bold = True
        p.add_run("，ROI 为 ")
        p.add_run(f"{summary.roi:.2f}").bold = True
        p.add_run("。")

        p2 = doc.add_paragraph()
        p2.add_run(f"期间共完成订单 {summary.order_count:,} 笔，客单价 ¥{summary.avg_order_value:.2f}，"
                  f"推广费率 {summary.promotion_rate*100:.2f}%，退款率 {summary.refund_rate*100:.2f}%。")

        # 说明
        doc.add_heading("三、说明", level=1)
        doc.add_paragraph("本报告由电商经营驾驶舱 Pro 自动生成。")

        file_path = self._generate_filename(f"{period}_report", "docx", store_id)
        doc.save(file_path)
        return file_path

    # ---------- PPT ----------
    def export_ppt(self, store_id: int = 0) -> Path:
        """导出经营总结 PPT（老板汇报神器）"""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        today = self.analytics.get_today_summary(store_id)
        week = self.analytics.get_week_summary(store_id)
        month = self.analytics.get_month_summary(store_id)
        year = self.analytics.get_year_summary(store_id)
        trend = self.analytics.get_trend(store_id, 30)
        sku_stats = self.analytics.get_sku_stats(store_id, 30)
        store_name = self._get_store_name(store_id)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 封面
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        # 背景色
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF7)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{store_name}"
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1D, 0x1D, 0x1F)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = f"经营总结报告 · {date.today().isoformat()}"
        run2.font.size = Pt(20)
        run2.font.color.rgb = RGBColor(0x6E, 0x6E, 0x73)

        # 销售概览
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title.text_frame.text = "一、销售概览"
        title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].runs[0].font.bold = True
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x71, 0xE3)

        # 四个数据卡片
        cards = [
            ("今日销售额", f"¥{today.sales_amount:,.0f}", f"{today.order_count} 单"),
            ("本周销售额", f"¥{week.sales_amount:,.0f}", f"{week.order_count} 单"),
            ("本月销售额", f"¥{month.sales_amount:,.0f}", f"{month.order_count} 单"),
            ("本年销售额", f"¥{year.sales_amount:,.0f}", f"{year.order_count} 单"),
        ]
        for idx, (label, value, sub) in enumerate(cards):
            x = Inches(0.5 + idx * 3.1)
            y = Inches(1.5)
            box = slide.shapes.add_textbox(x, y, Inches(3), Inches(2))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf = box.text_frame
            tf.text = label
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.size = Pt(14)
            p.runs[0].font.color.rgb = RGBColor(0x6E, 0x6E, 0x73)
            p2 = tf.add_paragraph()
            p2.text = value
            p2.alignment = PP_ALIGN.CENTER
            p2.runs[0].font.size = Pt(32)
            p2.runs[0].font.bold = True
            p2.runs[0].font.color.rgb = RGBColor(0x1D, 0x1D, 0x1F)
            p3 = tf.add_paragraph()
            p3.text = sub
            p3.alignment = PP_ALIGN.CENTER
            p3.runs[0].font.size = Pt(12)
            p3.runs[0].font.color.rgb = RGBColor(0x6E, 0x6E, 0x73)

        # 利润概览
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title.text_frame.text = "二、利润概览"
        title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].runs[0].font.bold = True
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x71, 0xE3)

        profit_cards = [
            ("今日净利润", f"¥{today.net_profit:,.0f}", f"利润率 {today.profit_rate*100:.1f}%"),
            ("本周净利润", f"¥{week.net_profit:,.0f}", f"利润率 {week.profit_rate*100:.1f}%"),
            ("本月净利润", f"¥{month.net_profit:,.0f}", f"利润率 {month.profit_rate*100:.1f}%"),
            ("本年净利润", f"¥{year.net_profit:,.0f}", f"利润率 {year.profit_rate*100:.1f}%"),
        ]
        for idx, (label, value, sub) in enumerate(profit_cards):
            x = Inches(0.5 + idx * 3.1)
            y = Inches(1.5)
            box = slide.shapes.add_textbox(x, y, Inches(3), Inches(2))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf = box.text_frame
            tf.text = label
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.size = Pt(14)
            p.runs[0].font.color.rgb = RGBColor(0x6E, 0x6E, 0x73)
            p2 = tf.add_paragraph()
            p2.text = value
            p2.alignment = PP_ALIGN.CENTER
            p2.runs[0].font.size = Pt(32)
            p2.runs[0].font.bold = True
            p2.runs[0].font.color.rgb = RGBColor(0x34, 0xC7, 0x59)  # Apple Green
            p3 = tf.add_paragraph()
            p3.text = sub
            p3.alignment = PP_ALIGN.CENTER
            p3.runs[0].font.size = Pt(12)
            p3.runs[0].font.color.rgb = RGBColor(0x6E, 0x6E, 0x73)

        # 趋势分析
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title.text_frame.text = "三、销售趋势（近 30 天）"
        title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].runs[0].font.bold = True
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x71, 0xE3)

        # 取近 30 天趋势数据
        if trend:
            avg_sales = sum(p.sales for p in trend) / len(trend)
            max_sales = max(p.sales for p in trend)
            min_sales = min(p.sales for p in trend)
            total_sales = sum(p.sales for p in trend)

            info_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
            tf = info_box.text_frame
            tf.text = f"近 30 天总销售额：¥{total_sales:,.2f}"
            p = tf.paragraphs[0]
            p.runs[0].font.size = Pt(18)
            p.runs[0].font.bold = True
            for line_text in [
                f"日均销售额：¥{avg_sales:,.2f}",
                f"最高单日：¥{max_sales:,.2f}",
                f"最低单日：¥{min_sales:,.2f}",
                f"近 7 天平均：¥{sum(p.sales for p in trend[-7:])/7:,.2f}",
            ]:
                p = tf.add_paragraph()
                p.text = line_text
                p.runs[0].font.size = Pt(16)

        # 爆款分析
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title.text_frame.text = "四、爆款分析 TOP 5"
        title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].runs[0].font.bold = True
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x71, 0xE3)

        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        tf = info_box.text_frame
        for idx, s in enumerate(sku_stats[:5], 1):
            p = tf.paragraphs[0] if idx == 1 else tf.add_paragraph()
            p.text = f"{idx}. {s.sku_name}（{s.sku_code}）"
            p.runs[0].font.size = Pt(20)
            p.runs[0].font.bold = True
            p2 = tf.add_paragraph()
            p2.text = f"    销售额 ¥{s.sales_amount:,.0f} · 销量 {s.quantity} 件 · 利润 ¥{s.gross_profit:,.0f} · ROI {s.roi:.2f}"
            p2.runs[0].font.size = Pt(14)
            p2.runs[0].font.color.rgb = RGBColor(0x6E, 0x6E, 0x73)

        # 问题与计划
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title.text_frame.text = "五、问题分析与下阶段计划"
        title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].runs[0].font.bold = True
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x71, 0xE3)

        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        tf = info_box.text_frame
        problems = [
            f"当前推广费率 {month.promotion_rate*100:.1f}%，建议控制在 20% 以内",
            f"退款率 {month.refund_rate*100:.1f}%，关注产品质量与售后",
            f"ROI {month.roi:.2f}，可优化推广渠道结构",
        ]
        for idx, text in enumerate(problems):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"• {text}"
            p.runs[0].font.size = Pt(16)

        p = tf.add_paragraph()
        p.text = "\n下阶段计划："
        p.runs[0].font.size = Pt(18)
        p.runs[0].font.bold = True

        plans = [
            "持续优化 TOP SKU 推广预算分配",
            "重点优化退款率高的 SKU",
            "拓展新爆款，分散销售集中度",
        ]
        for plan in plans:
            p = tf.add_paragraph()
            p.text = f"  ✓ {plan}"
            p.runs[0].font.size = Pt(14)

        file_path = self._generate_filename("business_ppt", "pptx", store_id)
        prs.save(file_path)
        return file_path

    def _get_store_name(self, store_id: int) -> str:
        if not store_id:
            return "全店铺汇总"
        store = self.session.get(Store, store_id)
        return store.name if store else "未知店铺"
